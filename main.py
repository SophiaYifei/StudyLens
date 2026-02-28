import os, re, json, warnings
warnings.filterwarnings("ignore")
from pathlib import Path
from typing import List, Dict, Any, Callable, Tuple

# File parsing
from pptx import Presentation                        # pip install python-pptx to read pptx files

# NLP library for text processing
import nltk
# split content into sentences
nltk.download("punkt_tab", quiet=True)
#  turns a paragraph (string) into a list of sentences.
from nltk.tokenize import sent_tokenize
from transformers import AutoTokenizer

# Embeddings
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt


#  Set Chunking parameters
CHUNK_SIZE_TOKENS   = 200      # max tokens per chunk
CHUNK_OVERLAP_SENTS = 1       # sentences of overlap between consecutive chunks
MIN_CHUNK_TOKENS    = 30      # discard very short chunks (empty slides, lone headers)

# Models
TOKENIZER_MODEL = "facebook/bart-large-cnn"
EMBEDDING_MODEL = "all-MiniLM-L6-v2"

_SCRIPT_DIR = Path(__file__).resolve().parent
DATA_DIR   = _SCRIPT_DIR / "data" / "raw"
OUTPUT_DIR = _SCRIPT_DIR / "outputs"


# Text cleaning helpers

def clean_transcript(raw: str) -> str:
    """
    Clean auto-generated lecture transcript:
      • Remove [Auto-generated …] header
      • Join short caption fragments into continuous prose
      • Collapse extra whitespace
    """
    lines = raw.splitlines()
    if lines and lines[0].startswith("[Auto-generated"):
        lines = lines[1:]
    text = " ".join(ln.strip() for ln in lines if ln.strip())
    return re.sub(r" {2,}", " ", text).strip()


def clean_notes(raw: str) -> str:
    """
    Light clean for student notes:
      • Normalize line endings
      • Collapse 3+ blank lines → 2
      • Preserve dividers and >>>> adversarial markers (test data!)
    """
    text = raw.replace("\r\n", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# Loaders

def load_txt(path: Path) -> Dict:
    raw = path.read_text(encoding="utf-8", errors="replace")
    if "Captions" in path.name:
        doc_type = "transcript"
        text = clean_transcript(raw)
    else:
        doc_type = "notes"
        text = clean_notes(raw)
    return {
        "source":   path.name,
        "doc_type": doc_type,
        "text":     text,
        "metadata": {"path": str(path), "chars": len(text)},
    }


def load_pptx(path: Path) -> Dict:
    """
    Extract text from every slide in a .pptx file.
    Title detected via placeholder_format.idx == 0.
    Returns full-doc text AND a per-slide list for structural chunking.
    """
    prs = Presentation(path)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text, body_parts = "", []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            frame = shape.text_frame.text.strip()
            if not frame:
                continue
            try:
                is_title = shape.is_placeholder and shape.placeholder_format.idx == 0
            except AttributeError:
                is_title = False
            if is_title:
                title_text = frame
            else:
                body_parts.append(frame)

        full = "\n".join(filter(None, [title_text] + body_parts)).strip()
        if full:
            slides_data.append({"slide_num": slide_num, "title": title_text, "text": full})

    full_doc = "\n\n--- SLIDE BREAK ---\n\n".join(s["text"] for s in slides_data)
    return {
        "source":   path.name,
        "doc_type": "slides",
        "text":     full_doc,
        "slides":   slides_data,
        "metadata": {"path": str(path), "num_slides": len(slides_data), "chars": len(full_doc)},
    }


def load_all_documents(data_dir: Path) -> List[Dict]:
    # Verify data loading
    documents = []
    for path in sorted(data_dir.iterdir()):
        ext = path.suffix.lower()
        if ext == ".txt":
            doc = load_txt(path)
            documents.append(doc)
            print(f"  [TXT / {doc['doc_type']:10s}]  {path.name[:58]}  ({doc['metadata']['chars']:>7,} chars)")
        elif ext == ".pptx":
            doc = load_pptx(path)
            documents.append(doc)
            print(f"  [PPTX/ slides    ]  {path.name[:58]}  ({doc['metadata']['num_slides']:>3} slides)")
        else:
            print(f"   [SKIP]  {path.name}")
    print(f"\n Loaded {len(documents)} documents total.")
    return documents


def verify_text_loaded(documents: List[Dict]) -> None:
    # Verify text loaded
    for doc in documents:
        label = (doc['metadata'].get('num_slides') and f"{doc['metadata']['num_slides']} slides") \
                or f"{doc['metadata']['chars']:,} chars"
        print(f"\n{'─'*72}")
        print(f"  source   : {doc['source']}")
        print(f"  doc_type : {doc['doc_type']}   ({label})")
        print(f"  preview  : {doc['text'][:220].replace(chr(10), ' ')} …")


def run_tokenizer_sanity_check(tokenizer) -> None:
    def count_tokens(text: str) -> int:
        """Count BPE tokens using the BART tokenizer (no special tokens)."""
        return len(tokenizer.encode(text, add_special_tokens=False))
    # Sanity check
    samples = [
        "Attention is a mechanism that computes a weighted sum of value vectors.",
        "self attn = contextual embeds via dot prod + softmax + weighted sum lol",
        "TF-IDF = TF * log(N / df_t)  where N = total docs, df_t = docs containing term",
    ]
    for s in samples:
        print(f"  {count_tokens(s):>3} tok  |  {s}")


# ── Strategy 1: Sentence-aware sliding window ─────────────────────────────────
# Best for: transcripts (continuous prose, no natural section markers)
#
# Algorithm:
#   1. Split text into sentences (NLTK)
#   2. Greedily accumulate sentences until adding the next would exceed chunk_size
#   3. Flush chunk; carry last `overlap_sents` sentences into the next chunk
#      (the triggering sentence is re-evaluated — nothing is skipped)
# Edge cases:
#   • Single sentence > chunk_size  → emitted as-is (can't split mid-sentence)
#   • Overlap alone   > chunk_size  → overlap discarded, fresh start

def chunk_by_sentences(
    text: str,
    count_tokens_fn,
    chunk_size: int = CHUNK_SIZE_TOKENS,
    overlap_sents: int = CHUNK_OVERLAP_SENTS,
    min_tokens: int = MIN_CHUNK_TOKENS,
) -> List[str]:
    sentences = sent_tokenize(text)
    if not sentences:
        # No sentence boundaries from NLTK (e.g. punkt_tab not loaded) - treat whole text as one chunk
        text_stripped = text.strip()
        if text_stripped and count_tokens_fn(text_stripped) >= min_tokens:
            return [text_stripped]
        return []

    chunks: List[str] = []
    cur_sents: List[str] = []
    cur_tokens: int = 0
    i = 0

    while i < len(sentences):
        sent = sentences[i]
        sent_tok = count_tokens_fn(sent)

        # Edge case: single sentence exceeds limit — emit as its own chunk
        if sent_tok > chunk_size and not cur_sents:
            chunks.append(sent)
            i += 1
            continue

        # Adding this sentence would overflow → flush current chunk
        if cur_tokens + sent_tok > chunk_size and cur_sents:
            chunk_text = " ".join(cur_sents)
            if count_tokens_fn(chunk_text) >= min_tokens:
                chunks.append(chunk_text)

            # Overlap: carry last N sentences into next chunk
            if overlap_sents and len(cur_sents) > overlap_sents:
                ov = cur_sents[-overlap_sents:]
                ov_tok = sum(count_tokens_fn(s) for s in ov)
                if ov_tok < chunk_size:       # overlap must itself fit
                    cur_sents, cur_tokens = ov, ov_tok
                else:
                    cur_sents, cur_tokens = [], 0
            else:
                cur_sents, cur_tokens = [], 0
            continue   # ← i NOT incremented: re-evaluate current sentence in new chunk

        cur_sents.append(sent)
        cur_tokens += sent_tok
        i += 1

    # Flush remainder
    if cur_sents:
        chunk_text = " ".join(cur_sents)
        if count_tokens_fn(chunk_text) >= min_tokens:
            chunks.append(chunk_text)

    return chunks


# ── Strategy 2: Structural (section-based) ────────────────────────────────────
# Best for: student notes (section dividers encode topic boundaries)
# Fallback: if a section > chunk_size, apply sentence-sliding-window within it

_SECTION_RE = re.compile(r"={5,}")


def chunk_by_structure(
    text: str,
    count_tokens_fn,
    chunk_size: int = CHUNK_SIZE_TOKENS,
    min_tokens: int = MIN_CHUNK_TOKENS,
) -> List[str]:
    sections = _SECTION_RE.split(text)
    chunks: List[str] = []

    for section in sections:
        section = section.strip()
        if not section:
            continue
        tok = count_tokens_fn(section)
        if tok <= chunk_size:
            if tok >= min_tokens:
                chunks.append(section)       # section fits as one chunk
        else:
            # Section too long → sentence-split it
            chunks.extend(chunk_by_sentences(section, count_tokens_fn, chunk_size, min_tokens=min_tokens))

    return chunks


# ── Strategy 3: Per-slide ─────────────────────────────────────────────────────
# Best for: PPTX slides (each slide = self-contained semantic unit)
# Short adjacent slides are merged into one chunk to avoid tiny chunks.
# Unusually long slides (dense bullet content) are sub-chunked via sentence sliding window.

def chunk_by_slides(
    slides: List[Dict],
    count_tokens_fn,
    chunk_size: int = CHUNK_SIZE_TOKENS,
    min_tokens: int = MIN_CHUNK_TOKENS,
) -> List[str]:
    chunks: List[str] = []
    buf_text, buf_tokens = "", 0

    for slide in slides:
        slide_text = slide["text"]
        slide_tok  = count_tokens_fn(slide_text)

        if slide_tok > chunk_size:
            # Flush buffer, then sub-chunk this long slide
            if buf_text and buf_tokens >= min_tokens:
                chunks.append(buf_text.strip())
            chunks.extend(chunk_by_sentences(slide_text, count_tokens_fn, chunk_size, min_tokens=min_tokens))
            buf_text, buf_tokens = "", 0
            continue

        if buf_tokens + slide_tok > chunk_size and buf_text:
            # Buffer full → flush, start new buffer with this slide
            if buf_tokens >= min_tokens:
                chunks.append(buf_text.strip())
            buf_text, buf_tokens = slide_text, slide_tok
        else:
            # Merge slide into buffer
            buf_text   = (buf_text + "\n\n" + slide_text).strip()
            buf_tokens += slide_tok

    if buf_text and buf_tokens >= min_tokens:
        chunks.append(buf_text.strip())

    return chunks


# ── Apply chunking to all documents ───────────────────────────────────────────

def apply_chunking_to_documents(documents: List[Dict], tokenizer) -> List[Dict]:
    def count_tokens(text: str) -> int:
        return len(tokenizer.encode(text, add_special_tokens=False))

    STRATEGY_MAP = {
        "transcript": ("sentence_sliding_window", lambda doc: chunk_by_sentences(doc["text"], count_tokens)),
        "notes":      ("structural_section",      lambda doc: chunk_by_structure(doc["text"], count_tokens)),
        "slides":     ("per_slide",               lambda doc: chunk_by_slides(doc["slides"], count_tokens)),
    }

    all_chunks: List[Dict] = []

    print(f"  {'Source':<52}  {'Type':12}  {'#Chunks':>7}  Strategy")
    print("  " + "─" * 96)

    for doc in documents:
        strategy_name, chunker_fn = STRATEGY_MAP[doc["doc_type"]]
        raw_chunks = chunker_fn(doc)

        stem = Path(doc["source"]).stem[:20]
        for idx, chunk_text in enumerate(raw_chunks):
            all_chunks.append({
                "chunk_id":    f"{stem}__c{idx:03d}",
                "source":      doc["source"],
                "doc_type":    doc["doc_type"],
                "chunk_index": idx,
                "strategy":    strategy_name,
                "text":        chunk_text,
                "token_count": count_tokens(chunk_text),
            })

        print(f"  {doc['source'][:50]:<50}  {doc['doc_type']:12}  {len(raw_chunks):>7}  {strategy_name}")

    print(f"\n  {'─'*96}")
    print(f"  Total chunks: {len(all_chunks)}")
    return all_chunks


# ── Chunk statistics + distribution plot ──────────────────────────────────────

def print_stats_and_plot(all_chunks: List[Dict]) -> Any:
    df = pd.DataFrame(all_chunks)

    print("=== Token count statistics by document type ===\n")
    stats = (
        df.groupby("doc_type")["token_count"]
        .agg(
            n_chunks="count",
            mean="mean",
            median="median",
            min="min",
            max="max",
            over_200=lambda x: (x > 200).sum(),
            over_256=lambda x: (x > 256).sum(),   # MiniLM hard limit
        )
        .round(1)
    )
    print(stats.to_string())

    # Flag any chunks that exceed MiniLM's limit
    over_256 = df[df["token_count"] > 256]
    if len(over_256):
        print(f"\n  {len(over_256)} chunk(s) exceed MiniLM's 256-token limit (will be truncated):")
        print(over_256[["chunk_id", "doc_type", "token_count"]].to_string(index=False))
    else:
        print(f"\n  All {len(df)} chunks fit within MiniLM's 256-token embedding limit.")

    # ── Plot ──────────────────────────────────────────────────────────────────────
    palette = {"transcript": "#55a868", "notes": "#c44e52", "slides": "#4c72b0"}
    fig, axes = plt.subplots(1, 2, figsize=(13, 4))

    # Left: overall distribution
    df["token_count"].hist(bins=30, ax=axes[0], color="#4C72B0", edgecolor="white", linewidth=0.5)
    axes[0].axvline(256, color="crimson",    linestyle="--", linewidth=1.8, label="MiniLM limit (256)")
    axes[0].axvline(200, color="darkorange", linestyle="--", linewidth=1.8, label=f"Target ({CHUNK_SIZE_TOKENS})")
    axes[0].set_title("Token Count — All Chunks", fontsize=12)
    axes[0].set_xlabel("Tokens")
    axes[0].set_ylabel("Count")
    axes[0].legend()

    # Right: by doc type
    for dtype, grp in df.groupby("doc_type"):
        grp["token_count"].hist(bins=20, ax=axes[1], alpha=0.65,
                             color=palette.get(dtype, "gray"), label=dtype)
    axes[1].axvline(256, color="crimson", linestyle="--", linewidth=1.8, label="MiniLM (256)")
    axes[1].set_title("Token Count by Document Type", fontsize=12)
    axes[1].set_xlabel("Tokens")
    axes[1].legend()

    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "chunk_distribution.png", dpi=130, bbox_inches="tight")
    plt.close()
    print(f" Saved {OUTPUT_DIR / 'chunk_distribution.png'}")
    return df


# ── Inspect sample chunks per type ────────────────────────────────────────────

def inspect_sample_chunks(df) -> None:
    for dtype in df["doc_type"].unique():
        sample = df[df["doc_type"] == dtype].iloc[0]
        print(f"\n  doc_type : {dtype}")
        print(f"  chunk_id : {sample['chunk_id']}")
        print(f"  tokens   : {sample['token_count']}")
        print(f"  strategy : {sample['strategy']}")
        print(f"  text     :\n{sample['text'][:400]}")
        print("  …")


# ── Sanity-check: semantic retrieval demo ─────────────────────────────────────
# If the embeddings are working correctly, topic-matched chunks should rank first.

def retrieve(query: str, embed_model, embeddings, all_chunks: List[Dict], top_k: int = 3) -> Any:
    """Embed query and return top-k most similar chunks via cosine similarity."""
    q_vec  = embed_model.encode([query], normalize_embeddings=True)
    scores = cosine_similarity(q_vec, embeddings)[0]
    top_idx = np.argsort(scores)[::-1][:top_k]
    return pd.DataFrame([{
        "rank":     r + 1,
        "score":    round(float(scores[i]), 4),
        "doc_type": all_chunks[i]["doc_type"],
        "source":   Path(all_chunks[i]["source"]).name[:38],
        "preview":  all_chunks[i]["text"][:100].replace("\n", " ") + "…",
    } for r, i in enumerate(top_idx)])


def main() -> None:
    OUTPUT_DIR.mkdir(exist_ok=True)

    all_files = sorted(DATA_DIR.iterdir())
    print(f"DATA_DIR  : {DATA_DIR.resolve()}")
    print(f"OUTPUT_DIR: {OUTPUT_DIR.resolve()}\n")
    print(f"  {'File':<68}  {'Ext':>5}  {'Size':>8}")
    print("  " + "─" * 84)
    for f in all_files:
        print(f"  {f.name[:68]:<68}  {f.suffix.upper():>5}  {f.stat().st_size/1024:>6.1f} KB")

    documents = load_all_documents(DATA_DIR)
    verify_text_loaded(documents)

    print(f"Loading tokenizer: {TOKENIZER_MODEL} …")
    tokenizer = AutoTokenizer.from_pretrained(TOKENIZER_MODEL)
    print(f" Tokenizer ready. Vocab size: {tokenizer.vocab_size:,}\n")
    run_tokenizer_sanity_check(tokenizer)

    all_chunks = apply_chunking_to_documents(documents, tokenizer)
    df = print_stats_and_plot(all_chunks)
    inspect_sample_chunks(df)

    print(f"Loading embedding model: {EMBEDDING_MODEL} …")
    embed_model = SentenceTransformer(EMBEDDING_MODEL)
    print(f" Model loaded")
    print(f"  Max sequence length : {embed_model.max_seq_length} tokens")
    print(f"  Embedding dimension : {embed_model.get_sentence_embedding_dimension()}")

    chunk_texts = [c["text"] for c in all_chunks]
    print(f"\nEmbedding {len(chunk_texts)} chunks …")
    embeddings = embed_model.encode(
        chunk_texts,
        batch_size=32,
        show_progress_bar=True,
        normalize_embeddings=True,    # L2-normalise → cosine sim = dot product
        convert_to_numpy=True,
    )
    print(f"\n Embeddings: {embeddings.shape}   ({embeddings.shape[0]} chunks × {embeddings.shape[1]} dims)")
    print(f"  dtype  : {embeddings.dtype}")
    print(f"  norm[0]: {np.linalg.norm(embeddings[0]):.5f}   (≈ 1.0 after L2 normalisation)")

    test_queries = [
        "Neural Networks for NLP: word embeddings word2vec GloVe RNN LSTM GRU sequence models",
        "Attention & Transformers: self-attention keys queries values multi-head transformer BERT GPT",
       # "self-attention mechanism and how it produces contextual embeddings",
        #"text preprocessing steps: tokenization stemming lemmatization",
        #"ROUGE metric for evaluating summarization quality",
        #"BERT vs GPT architecture differences",
    ]
    for q in test_queries:
        print(f"\n  QUERY: {q}")
        print(retrieve(q, embed_model, embeddings, all_chunks).to_string(index=False))

    # 1. chunks.json  (text + metadata, no embedding arrays)
    chunks_path = OUTPUT_DIR / "chunks.json"
    with open(chunks_path, "w", encoding="utf-8") as f:
        json.dump(all_chunks, f, indent=2, ensure_ascii=False)
    print(f" chunks.json : {chunks_path}")

    # 2. embeddings.npy  (float32, L2-normalised)
    embed_path = OUTPUT_DIR / "embeddings.npy"
    np.save(embed_path, embeddings)
    print(f" embeddings.npy : {embed_path}   shape={embeddings.shape}")

    # 3. chunk_metadata.csv
    csv_path = OUTPUT_DIR / "chunk_metadata.csv"
    df.to_csv(csv_path, index=False)
    print(f" chunk_metadata.csv : {csv_path}")

    # 4. Reload verification
    loaded_chunks = json.load(open(chunks_path, encoding="utf-8"))
    loaded_embeds = np.load(embed_path)
    assert len(loaded_chunks) == len(loaded_embeds), "Length mismatch!"
    assert loaded_embeds.shape[1] == 384, "Wrong embedding dim!"
    print(f"\n Reload check passed: {len(loaded_chunks)} chunks, embeddings {loaded_embeds.shape}")

    print(f"""

  Documents loaded      : {len(documents):<35}
  Total chunks          : {len(all_chunks):<35}
  Chunks > 256 tok      : {int((df['token_count'] > 256).sum()):<35}
  Embedding shape       : {str(embeddings.shape):<35}
  Embedding model       : {EMBEDDING_MODEL:<35}
  Outputs saved to      : outputs/      """)


if __name__ == "__main__":
    main()

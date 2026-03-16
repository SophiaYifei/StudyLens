"""
app.py - StudyLens Web Application Backend

AI-powered lecture review platform with real-time summarization and RAG Q&A.
Data files are stored in a private HuggingFace dataset repo and downloaded
on first startup.

Usage:
    export OPENROUTER_API_KEY=your-key
    export HF_TOKEN=your-hf-token
    python app.py

Attribution: FastAPI framework, OpenRouter API, HuggingFace Hub.
"""

import os
import logging
from pathlib import Path
from fastapi import FastAPI, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pydantic import BaseModel

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("studylens")

# ── Config ───────────────────────────────────────────────────────────────
ROOT_DIR = Path(__file__).resolve().parent
DATA_DIR = ROOT_DIR / "data"
DATA_RAW = DATA_DIR / "raw"
DATA_PROCESSED = DATA_DIR / "processed"
DATA_OUTPUTS = DATA_DIR / "outputs"
USER_NOTES_DIR = DATA_DIR / "user_notes"
USER_NOTES_DIR.mkdir(parents=True, exist_ok=True)

OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY", "")
HF_TOKEN = os.environ.get("HF_TOKEN", "")
HF_DATASET_REPO = "YifeiGuo/studylens-data"  # Change to your repo

# ── Model Selection ──────────────────────────────────────────────────────
# Uncomment the one you want:

# FREE - good quality, moderate speed
# OPENROUTER_MODEL = "meta-llama/llama-3.3-70b-instruct:free"

# FREE - strong but slow
# OPENROUTER_MODEL = "nvidia/nemotron-3-super-120b-a12b:free"

# PAID ~$0.25/1M input - best value, fast, excellent quality
OPENROUTER_MODEL = "deepseek/deepseek-chat"

# PAID ~$0.10/1M input - cheapest and fastest
# OPENROUTER_MODEL = "google/gemini-2.0-flash-001"


# ── HuggingFace Data Download ────────────────────────────────────────────

def download_data():
    """Download data from private HuggingFace dataset repo on first run."""
    marker = DATA_DIR / ".hf_downloaded"
    if marker.exists():
        logger.info("Data already downloaded, skipping.")
        return

    if not HF_TOKEN:
        logger.warning("HF_TOKEN not set. Skipping data download. "
                       "Data must already exist locally.")
        return

    logger.info(f"Downloading data from {HF_DATASET_REPO}...")
    try:
        from huggingface_hub import snapshot_download
        snapshot_download(
            repo_id=HF_DATASET_REPO,
            repo_type="dataset",
            token=HF_TOKEN,
            local_dir=str(DATA_DIR),
            local_dir_use_symlinks=False,
        )
        marker.touch()
        logger.info("Data download complete.")
    except Exception as e:
        logger.error(f"Failed to download data: {e}")
        logger.info("Continuing with whatever data exists locally.")


# ── Topic metadata ───────────────────────────────────────────────────────

TOPICS = [
    {"id": "dl_s1", "course": "Deep Learning", "week": 1,
     "title": "Introduction to Deep Learning"},
    {"id": "dl_s2", "course": "Deep Learning", "week": 2,
     "title": "Computer Vision I"},
    {"id": "dl_s3", "course": "Deep Learning", "week": 3,
     "title": "Computer Vision II & Generative Models"},
    {"id": "dl_s4", "course": "Deep Learning", "week": 6,
     "title": "Natural Language Processing I"},
    {"id": "dl_s5", "course": "Deep Learning", "week": 7,
     "title": "Natural Language Processing II"},
    {"id": "ml_s1", "course": "Machine Learning", "week": 1,
     "title": "Introduction to Machine Learning"},
    {"id": "ml_s2", "course": "Machine Learning", "week": 2,
     "title": "Supervised Learning"},
    {"id": "ml_s3", "course": "Machine Learning", "week": 3,
     "title": "Unsupervised Learning"},
    {"id": "ml_s4", "course": "Machine Learning", "week": 4,
     "title": "Evaluation & Model Selection"},
    {"id": "ml_s5", "course": "Machine Learning", "week": 5,
     "title": "Advanced Topics in ML"},
]


def _get_topic(topic_id: str) -> dict:
    """Look up topic metadata by ID."""
    for t in TOPICS:
        if t["id"] == topic_id:
            return t
    raise HTTPException(status_code=404, detail=f"Topic {topic_id} not found")


def _read_file(path: Path) -> str:
    """Read a text file, returning empty string if not found."""
    return path.read_text(encoding="utf-8", errors="replace") if path.exists() else ""


# ── Data accessors ───────────────────────────────────────────────────────

def _extract_slide_text(topic_id: str) -> str:
    """Extract text from pptx for AI context (not for display)."""
    pptx_path = DATA_RAW / f"{topic_id}.pptx"
    if pptx_path.exists():
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_path))
            texts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        t = shape.text.strip()
                        if t in ("<#>", "\u2039#\u203a", "#") or len(t) < 2:
                            continue
                        slide_texts.append(t)
                if slide_texts:
                    texts.append(
                        f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))
            return "\n\n".join(texts)
        except Exception:
            pass
    return _read_file(DATA_PROCESSED / f"{topic_id}_ori.txt")


def _get_transcript(topic_id: str) -> str:
    """Read cleaned transcript for a topic."""
    num = topic_id.split("_s")[1]
    prefix = topic_id.split("_")[0]
    cleaned = DATA_RAW / f"{prefix}_t{num}_cleaned.txt"
    if cleaned.exists():
        return _read_file(cleaned)
    return _read_file(DATA_RAW / "transcripts" / f"{prefix}_t{num}.txt")


def _get_notes(topic_id: str) -> str:
    """Read notes: user-saved first, then original."""
    user_path = USER_NOTES_DIR / f"{topic_id}_notes.txt"
    if user_path.exists():
        return _read_file(user_path)
    num = topic_id.split("_s")[1]
    prefix = topic_id.split("_")[0]
    return _read_file(DATA_RAW / f"{prefix}_n{num}.txt")


def _get_precomputed_summary(topic_id: str) -> str:
    """Read pre-computed Qwen-FT summary if available."""
    ft_path = (DATA_OUTPUTS / "finetune" / "qwen7b-ft" / "final"
               / f"{topic_id}_sum_qwen7b-ft_final.txt")
    return _read_file(ft_path)


# ── OpenRouter API ───────────────────────────────────────────────────────

async def _call_openrouter(
    system: str, messages: list, max_tokens: int = 4096
) -> str:
    """Call OpenRouter chat completions API."""
    if not OPENROUTER_API_KEY:
        raise HTTPException(
            status_code=500, detail="OPENROUTER_API_KEY not set")
    import httpx
    all_messages = [{"role": "system", "content": system}] + messages
    async with httpx.AsyncClient(timeout=180.0) as client:
        resp = await client.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": OPENROUTER_MODEL,
                "messages": all_messages,
                "max_tokens": max_tokens,
            },
        )
    if resp.status_code != 200:
        raise HTTPException(
            status_code=502, detail=f"OpenRouter error: {resp.text}")
    return resp.json()["choices"][0]["message"]["content"]


# ── FastAPI Application ──────────────────────────────────────────────────

app = FastAPI(title="StudyLens", version="1.0")


@app.on_event("startup")
async def startup():
    """Download data from HuggingFace on first startup."""
    download_data()


static_dir = ROOT_DIR / "static"
if static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")


@app.get("/")
async def index():
    """Serve the frontend."""
    return FileResponse(str(static_dir / "index.html"))


# ── API: Topics ──────────────────────────────────────────────────────────

@app.get("/api/topics")
async def list_topics():
    """List all available lecture topics."""
    return TOPICS


@app.get("/api/topics/{topic_id}")
async def get_topic_detail(topic_id: str):
    """Get full detail for a topic including slides text, transcript, notes."""
    topic = _get_topic(topic_id)
    pdf_exists = (DATA_RAW / f"{topic_id}.pdf").exists()
    return {
        **topic,
        "slides_text": _extract_slide_text(topic_id)[:50000],
        "transcript": _get_transcript(topic_id)[:50000],
        "notes": _get_notes(topic_id),
        "precomputed_summary": _get_precomputed_summary(topic_id),
        "has_slides_pdf": pdf_exists,
    }


# ── API: Slides PDF ─────────────────────────────────────────────────────

@app.get("/api/slides-pdf/{topic_id}")
async def get_slides_pdf(topic_id: str):
    """Serve PDF slides for a topic."""
    pdf_path = DATA_RAW / f"{topic_id}.pdf"
    if pdf_path.exists():
        return FileResponse(str(pdf_path), media_type="application/pdf")
    raise HTTPException(status_code=404, detail="PDF not found.")


# ── API: Notes ───────────────────────────────────────────────────────────

class NotesUpdate(BaseModel):
    topic_id: str
    content: str


@app.post("/api/notes")
async def save_notes(data: NotesUpdate):
    """Save user-edited notes for a topic."""
    path = USER_NOTES_DIR / f"{data.topic_id}_notes.txt"
    path.write_text(data.content, encoding="utf-8")
    return {"status": "saved"}


# ── API: Summarize ───────────────────────────────────────────────────────

class SummarizeRequest(BaseModel):
    topic_id: str
    notes_content: str = ""


@app.post("/api/summarize")
async def summarize(req: SummarizeRequest):
    """Generate an AI summary from slides + transcript + notes."""
    slides = _extract_slide_text(req.topic_id)
    transcript = _get_transcript(req.topic_id)
    combined = (f"SLIDE TEXT:\n{slides[:20000]}\n\n"
                f"TRANSCRIPT:\n{transcript[:20000]}")
    if req.notes_content:
        combined += f"\n\nSTUDENT NOTES:\n{req.notes_content[:5000]}"

    system = (
        "You are StudyLens, an expert academic summarizer for CS courses. "
        "Create clear, structured summaries for exam review. "
        "Use markdown: ## headers, **bold** key terms, - bullets, "
        "`code` for notation. "
        "For math, use LaTeX with $...$ inline and $$...$$ display. "
        "IMPORTANT: Always use $$ for display math, never \\[ or \\]. "
        "Keep all LaTeX on a single line, never split across lines."
    )
    user_msg = (
        "Create a comprehensive study summary. "
        "Ignore filler/off-topic content.\n\n"
        "Structure:\n"
        "## Key Concepts\n## Technical Details\n"
        "## Important Relationships\n## Practical Takeaways\n\n"
        "Stay under 1200 words. Complete every section.\n\n" + combined
    )
    summary = await _call_openrouter(
        system, [{"role": "user", "content": user_msg}])
    return {"summary": summary}


# ── API: Chat (RAG) ─────────────────────────────────────────────────────

class ChatRequest(BaseModel):
    topic_id: str
    question: str
    notes_content: str = ""
    history: list = []


@app.post("/api/chat")
async def chat(req: ChatRequest):
    slides = _extract_slide_text(req.topic_id)
    transcript = _get_transcript(req.topic_id)
    
    # ADD THIS LINE for debugging
    logger.info(f"Chat context: slides={len(slides)} chars, transcript={len(transcript)} chars, notes={len(req.notes_content)} chars")
    
    context = (f"SLIDE TEXT:\n{slides[:15000]}\n\n"
               f"TRANSCRIPT:\n{transcript[:15000]}")
    if req.notes_content:
        context += f"\n\nYOUR NOTES:\n{req.notes_content[:3000]}"

    system = (
        "You are StudyLens, the student's personal study assistant. "
        "You MUST answer based on the LECTURE MATERIALS provided below. "
        "ALWAYS cite specific slide numbers or transcript content. "
        "If a student asks a concept question, find the relevant part in the slides/transcript and reference it. "
        "If you cannot find the answer in the materials, say 'I couldn't find this in your lecture materials.'\n\n"
        "IMPORTANT: The lecture materials are provided after this system message. READ THEM CAREFULLY.\n\n"
        "Style: Use 'you/your', be warm like a TA. Use markdown and $...$ for math.\n\n"
        f"=== LECTURE MATERIALS ===\n{context}"
    )
    messages = [
        {"role": h["role"], "content": h["content"]}
        for h in req.history[-6:]
    ]
    messages.append({"role": "user", "content": req.question})
    answer = await _call_openrouter(system, messages, max_tokens=1500)
    return {"answer": answer}


# ── Run ──────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 7860))
    uvicorn.run(app, host="0.0.0.0", port=port)
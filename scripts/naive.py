"""
scripts/naive.py - Naive baseline: extract first sentence from first 5 slides.

Produces one _naive.txt per PPTX file in data/outputs/naive/.
The FirstSentenceSummarizer class in model.py provides the same logic
as a BaseSummarizer subclass for the evaluation pipeline.

AI Attribution: Code co-authored with Claude (Anthropic, https://claude.ai)
for structural design, debugging, and documentation.
"""

from pathlib import Path
from pptx import Presentation
from zipfile import BadZipFile
import re

# Default: input dir for PPTs when run as script; output written to data/outputs/naive
ROOT_DIR = Path(__file__).resolve().parent.parent
DEFAULT_OUTPUT_DIR = ROOT_DIR / "data" / "outputs" / "naive"


def get_first_sentence_from_first_5_slides(filepath):
    prs = Presentation(filepath)
    collected_sentences = []

    for i, slide in enumerate(prs.slides):
        if i >= 5:
            break

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()

                sentences = re.split(r'(?<=[.!?])\s+', text)

                for sentence in sentences:
                    clean = sentence.strip()
                    if clean:
                        collected_sentences.append(clean)
                        break  # first sentence from this slide

        # continue to next slide automatically

    return collected_sentences


def process_all_ppts(directory: Path, output_dir=None):
    output_dir = output_dir or DEFAULT_OUTPUT_DIR
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    for ppt_file in directory.glob("*.pptx"):
        if ppt_file.name.startswith("~$"):
            continue

        try:
            sentences = get_first_sentence_from_first_5_slides(ppt_file)
            out_path = output_dir / (ppt_file.stem + "_naive.txt")
            out_path.write_text("\n".join(sentences) + "\n", encoding="utf-8")
            print(f"Created: {out_path}")
        except BadZipFile:
            print(f"Skipped (not valid pptx): {ppt_file}")
        except Exception as e:
            print(f"Error processing {ppt_file}: {e}")


if __name__ == "__main__":
    # When run as script, use data/raw as input and write to data/outputs/naive
    input_dir = ROOT_DIR / "data" / "raw"
    process_all_ppts(input_dir, DEFAULT_OUTPUT_DIR)
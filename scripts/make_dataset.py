"""
scripts/make_dataset.py
Data loading and text cleaning for lecture materials (PPTX slides,
auto-generated transcripts, student notes).
"""

import re
from pathlib import Path
from typing import List, Dict

from pptx import Presentation


# ── Text cleaning ──────────────────────────────────────────────────────

def clean_transcript(raw: str) -> str:
    """Remove [Auto-generated ...] header, join caption fragments, collapse whitespace."""
    lines = raw.splitlines()
    if lines and lines[0].startswith("[Auto-generated"):
        lines = lines[1:]
    text = " ".join(ln.strip() for ln in lines if ln.strip())
    return re.sub(r" {2,}", " ", text).strip()


def clean_notes(raw: str) -> str:
    """Normalize line endings and collapse 3+ blank lines to 2."""
    text = raw.replace("\r\n", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ── Loaders ────────────────────────────────────────────────────────────

def load_txt(path: Path) -> Dict:
    raw = path.read_text(encoding="utf-8", errors="replace")
    if "Captions" in path.name:
        doc_type = "transcript"
        text = clean_transcript(raw)
    else:
        doc_type = "notes"
        text = clean_notes(raw)
    return {
        "source":   path.name,
        "doc_type": doc_type,
        "text":     text,
        "metadata": {"path": str(path), "chars": len(text)},
    }


def load_pptx(path: Path) -> Dict:
    """Extract text from every slide. Returns full-doc text and per-slide list."""
    prs = Presentation(path)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text, body_parts = "", []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            frame = shape.text_frame.text.strip()
            if not frame:
                continue
            try:
                is_title = shape.is_placeholder and shape.placeholder_format.idx == 0
            except AttributeError:
                is_title = False
            if is_title:
                title_text = frame
            else:
                body_parts.append(frame)

        full = "\n".join(filter(None, [title_text] + body_parts)).strip()
        if full:
            slides_data.append({"slide_num": slide_num, "title": title_text, "text": full})

    full_doc = "\n\n--- SLIDE BREAK ---\n\n".join(s["text"] for s in slides_data)
    return {
        "source":   path.name,
        "doc_type": "slides",
        "text":     full_doc,
        "slides":   slides_data,
        "metadata": {"path": str(path), "num_slides": len(slides_data), "chars": len(full_doc)},
    }


def load_all_documents(data_dir: Path) -> List[Dict]:
    """Load all .txt and .pptx files from data_dir."""
    documents = []
    for path in sorted(data_dir.iterdir()):
        ext = path.suffix.lower()
        if ext == ".txt":
            documents.append(load_txt(path))
        elif ext == ".pptx":
            documents.append(load_pptx(path))
    return documents


# Remove noise from transcripts
import re
from pathlib import Path


def denoise_all_transcripts(directory: Path):
    """
    - Finds all .txt files containing 't'
    - Denoises them
    - Writes *_cleaned.txt
    - Prints wordcount before/after
    """

    # ---------- CONFIG ----------
    FILLER_WORDS = [
        "um", "uh", "erm", "like", "you know", "i mean",
        "sort of", "kind of", "right", "okay", "ok",
        "alright", "so", "well", "yeah", "yep", "nope"
    ]

    NOISE_LINE_PATTERNS = [
        r"^\s*(thanks|thank you).*?$",
        r"^\s*(welcome|hi|hello).*?$",
        r"^\s*(can you hear me|audio check).*?$",
        r"^\s*(let('?s)? get started).*?$",
        r"^\s*(any questions\??|questions\??)\s*$",
        r"^\s*(break|short break|quick break).*?$",
        r"^\s*(welcome back|we('?re)? back).*?$",
        r"^\s*(wrap( ?)?up|in conclusion).*?$",
        r"^\s*(bye|goodbye|see you).*?$",
    ]

    # Remove question lines (often Q&A noise). 
    DROP_QUESTION_LINES = True

    # Trim these many lines at very start/end if they match noise patterns
    MAX_TRIM_START_LINES = 50
    MAX_TRIM_END_LINES = 50
    # --------------------------------


    def wordcount(path: Path) -> int:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return len(re.findall(r"\b\w+\b", text))


    def clean_text_file(input_path: Path, output_path: Path):
        raw = input_path.read_text(encoding="utf-8", errors="ignore")
        lines = raw.splitlines()

        compiled_noise = [re.compile(p, re.IGNORECASE) for p in NOISE_LINE_PATTERNS]

        # ---- Trim start/end ----
        start, end = 0, len(lines)

        for i in range(min(MAX_TRIM, len(lines))):
            if any(rx.match(lines[i].strip()) for rx in compiled_noise):
                start = i + 1
            else:
                break

        for i in range(len(lines) - 1, max(-1, len(lines) - MAX_TRIM - 1), -1):
            if any(rx.match(lines[i].strip()) for rx in compiled_noise):
                end = i
            else:
                break

        lines = lines[start:end]

        # ---- Drop noise lines ----
        cleaned_lines = []
        for ln in lines:
            raw_ln = ln.strip()

            if not raw_ln:
                continue

            if any(rx.match(raw_ln) for rx in compiled_noise):
                continue

            if DROP_QUESTION_LINES and ("?" in raw_ln):
                continue

            cleaned_lines.append(ln)

        text = "\n".join(cleaned_lines)

        # ---- Remove fillers ----
        for fw in sorted(FILLER_WORDS, key=len, reverse=True):
            text = re.sub(rf"\b{re.escape(fw)}\b", "", text, flags=re.IGNORECASE)

        # remove repeated words
        text = re.sub(r"\b(\w+)(\s+\1\b)+", r"\1", text, flags=re.IGNORECASE)

        # clean spacing
        text = re.sub(r"\s+([,.!?;:])", r"\1", text)
        text = re.sub(r"[ \t]{2,}", " ", text).strip()

        output_path.write_text(text + "\n", encoding="utf-8")


    # ---- Main Processing ----
    target_files = [
        f for f in directory.glob("*.txt")
        if "t" in f.stem.lower() and not f.stem.endswith("_cleaned")
    ]

    print(f"Found {len(target_files)} matching files.\n")

    for file_path in target_files:
        output_path = file_path.with_name(file_path.stem + "_cleaned.txt")

        print(f"Denoising: {file_path.name}")

        before = wordcount(file_path)
        clean_text_file(file_path, output_path)
        after = wordcount(output_path)

        removed = before - after
        pct = round((removed / before * 100), 2) if before else 0

        print(f"Before: {before}")
        print(f"After : {after}")
        print(f"Removed: {removed} ({pct}%)\n")



"""
scripts/model.py - Summarization models for StudyLens.

AI Attribution: Code co-authored with Claude (Anthropic, https://claude.ai)
for structural design, debugging, and documentation.
"""

import math
import os
import random
import re
from pathlib import Path

import nltk
import torch
from nltk.tokenize import sent_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer
from transformers import pipeline, AutoTokenizer

_PROJECT_ROOT = Path(__file__).resolve().parent.parent

# Load .env file if it exists (optional, for local development)
try:
    from dotenv import load_dotenv
    load_dotenv(_PROJECT_ROOT / ".env", override=True)
except ImportError:
    pass  # dotenv not installed, rely on system environment variables


nltk.download('punkt_tab')

DEVICE = 0 if torch.cuda.is_available() else -1


# --- Base Class ---
class BaseSummarizer:
    """All summarizers must implement .summarize(text, final_pass: bool = True) -> str"""

    def summarize(self, text: str, final_pass: bool = True) -> str:
        raise NotImplementedError("Subclasses must implement summarize()")

# --- Naive Baseline: First Sentence from First 5 Slides ---
class FirstSentenceSummarizer(BaseSummarizer):
    """
    Naive baseline: extract the first sentence from each of the first 5 slides.

    There are two usage modes:
      * `summarize(text)`: expects preprocessed slide text (as used in the pipeline)
        and returns the first sentence from up to the first `num_slides` sentences.
      * `summarize_from_pptx(filepath)`: expects a path to a `.pptx` file and
        extracts the first sentence from each of the first `num_slides` slides.

    This is a heuristic approach assuming slide titles/first sentences capture
    the main topic of each slide.
    """

    def __init__(self, num_slides=5):
        self.num_slides = num_slides

    def summarize_from_pptx(self, filepath: str) -> str:
        """Summarize a `.pptx` by extracting the first sentence from the first N slides."""
        from pptx import Presentation
        from zipfile import BadZipFile

        try:
            prs = Presentation(filepath)
        except BadZipFile as exc:
            # Provide a clear, deterministic error for invalid/corrupted .pptx inputs
            raise ValueError(f"Invalid or corrupted .pptx file: {filepath}") from exc

        collected_sentences = []

        for i, slide in enumerate(prs.slides):
            if i >= self.num_slides:
                break

            slide_sentence_found = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    sentences = re.split(r'(?<=[.!?])\s+', text)
                    for sentence in sentences:
                        clean = sentence.strip()
                        if clean:
                            collected_sentences.append(clean)
                            slide_sentence_found = True
                            break
                    if slide_sentence_found:
                        break
        return " ".join(collected_sentences)

    def summarize(self, text: str, final_pass: bool = True) -> str:
        """
        Summarize preprocessed slide text by taking the first sentence from up to
        `num_slides` sentences in the input.

        This is used as a fallback/text-only baseline when the original `.pptx`
        file is not available.
        """
        sentences = sent_tokenize(text)
        k = min(self.num_slides, len(sentences))
        return " ".join(sentences[:k])


# --- Naive Baseline: Random Extractive ---
class RandomExtractiveSummarizer(BaseSummarizer):
    """
    Naive baseline: randomly select sentences from the input text.
    No model, no learning — just random sampling.
    This establishes the absolute floor that any real model should beat.
    """

    def __init__(self, num_sentences=15, seed=42):
        """
        Args:
            num_sentences: how many sentences to randomly pick
            seed: random seed for reproducibility
        """
        self.num_sentences = num_sentences
        self.seed = seed

    def summarize(self, text: str, final_pass: bool = True) -> str:
        sentences = sent_tokenize(text)
        print(f"Input: {len(sentences)} sentences.")

        # Pick random sentences (or all if fewer than num_sentences)
        k = min(self.num_sentences, len(sentences))
        random.seed(self.seed)
        selected_indices = sorted(random.sample(range(len(sentences)), k))

        # Keep original order so the summary reads coherently
        summary = " ".join(sentences[i] for i in selected_indices)
        print(f"Randomly selected {k} sentences.")
        return summary


# --- Classical ML: TF-IDF Extractive ---
class TFIDFExtractiveSummarizer(BaseSummarizer):
    """
    Classical ML baseline: rank sentences by TF-IDF importance score,
    then select the top-k most informative sentences.

    How it works:
    1. Split input into sentences
    2. Fit TF-IDF on the sentences (each sentence = one "document")
    3. Score each sentence = sum of its TF-IDF values
       (sentences with rare, informative words get higher scores)
    4. Select top-k highest-scoring sentences
    5. Return them in original order

    This is a standard extractive summarization approach using
    a classical (non-neural) ML technique.
    """

    def __init__(self, num_sentences=15):
        """
        Args:
            num_sentences: how many top sentences to select
        """
        self.num_sentences = num_sentences

    def summarize(self, text: str, final_pass: bool = True) -> str:
        sentences = sent_tokenize(text)
        print(f"Input: {len(sentences)} sentences.")

        if len(sentences) <= self.num_sentences:
            return text

        # Fit TF-IDF on the sentences
        # Each sentence is treated as a "document" in the corpus
        vectorizer = TfidfVectorizer(stop_words="english")
        tfidf_matrix = vectorizer.fit_transform(sentences)

        # Score each sentence: sum of TF-IDF values across all words
        # Higher score = more unique/informative content
        sentence_scores = tfidf_matrix.sum(axis=1).A1  # .A1 converts matrix to flat array

        # Get indices of top-k scoring sentences
        k = min(self.num_sentences, len(sentences))
        top_indices = sentence_scores.argsort()[::-1][:k]

        # Sort by original position to maintain document flow
        top_indices = sorted(top_indices)

        summary = " ".join(sentences[i] for i in top_indices)
        print(f"Selected top {k} sentences by TF-IDF score.")
        return summary

# --- BART Summarizer ---
class BARTSummarizer(BaseSummarizer):

    def __init__(self):
        """Initialize BART-CNN summarization pipeline."""
        # Load the HuggingFace pipeline once
        self.model_name = "facebook/bart-large-cnn"
        self.pipe = pipeline("summarization", model=self.model_name, device=DEVICE)
        self.tokenizer = AutoTokenizer.from_pretrained(self.model_name)
        self.max_input_tokens = 1024

    def _count_tokens(self, text):
        """Count BPE tokens in text."""
        # Use self.tokenizer to count tokens in text
        # Return integer
        return len(self.tokenizer.encode(text, add_special_tokens=False))

    def _split_into_chunks(self, text: str) -> list[str]:
        """
        Greedily fill each chunk sentence by sentence.
        When adding the next sentence would exceed SAFE_LIMIT,
        close the current chunk and carry that sentence to the next one.
        Single sentences that exceed SAFE_LIMIT are hard-truncated at token level.
        """
        SAFE_LIMIT = self.max_input_tokens - 2   # reserve 2 for BOS/EOS

        sentences = sent_tokenize(text)
        chunks = []
        current_sentences = []
        current_tokens = 0

        for sentence in sentences:
            sent_toks = self._count_tokens(sentence)

            # Edge case: a single sentence is already over the limit
            # Hard-truncate it at token level and treat as its own chunk
            if sent_toks > SAFE_LIMIT:
                # First, flush whatever we have so far
                if current_sentences:
                    chunks.append(" ".join(current_sentences))
                    current_sentences = []
                    current_tokens = 0
                # Truncate the oversized sentence and save it directly
                token_ids = self.tokenizer.encode(sentence, add_special_tokens=False)
                truncated = self.tokenizer.decode(token_ids[:SAFE_LIMIT], skip_special_tokens=True)
                chunks.append(truncated)
                continue

            # Normal case: adding this sentence would exceed the limit
            # → close current chunk, start a new one with this sentence
            if current_tokens + sent_toks > SAFE_LIMIT:
                chunks.append(" ".join(current_sentences))
                current_sentences = []
                current_tokens = 0

            current_sentences.append(sentence)
            current_tokens += sent_toks

        # Don't forget the last chunk
        if current_sentences:
            chunks.append(" ".join(current_sentences))

        return chunks


    def _summarize_single(self, text: str, max_length=150, min_length=40) -> str:
        """Summarize a single chunk within token limits."""
        # Call self.pipe(text, max_length=..., min_length=..., do_sample=False)
        # Return the summary_text string

        # Failsafe: hard-truncate at token level before passing to model
        token_ids = self.tokenizer.encode(text, add_special_tokens=False)
        if len(token_ids) > self.max_input_tokens - 2:
            token_ids = token_ids[:self.max_input_tokens - 2]
            text = self.tokenizer.decode(token_ids, skip_special_tokens=True)

        # Safety checks
        text_tokens = self._count_tokens(text)
        if text_tokens < min_length:
            # Adjust min_length downwards for very short inputs
            min_length = max(10, text_tokens // 2)
        if max_length <= min_length:
            # Ensure max_length is always greater than min_length
            max_length = min_length + 20

        result = self.pipe(
            text,
            max_length=max_length,
            min_length=min_length,
            do_sample=False,
            truncation=True
        )
        return result[0]['summary_text']

    def summarize(self, text: str, final_pass: bool = True) -> str:         # This is the main method called from outside
        # Step 1: Count tokens of input text
        total_tokens = self._count_tokens(text)

        # Step 2: If tokens <= 1024, just call _summarize_single directly
        if total_tokens <= self.max_input_tokens:
            print(f"Text fits in one pass ({total_tokens} tokens), summarizing directly.")
            max_len = min(total_tokens // 2, 512)
            return self._summarize_single(text, max_length=max_len, min_length=40)
        
        # Step 3: If tokens > 1024 (hierarchical summarization):
        print(f"Text too long ({total_tokens} tokens), performing hierarchical summarization.")

        #   a. chunks = self._split_into_chunks(text)
        chunks = self._split_into_chunks(text)
        print(f"Split text into {len(chunks)} chunks.")

        #   b. chunk_summaries = []
        #      for each chunk:
        #          summary = self._summarize_single(chunk)
        #          chunk_summaries.append(summary)
        chunk_summaries = []
        for i, chunk in enumerate(chunks):
            chunk_tokens = self._count_tokens(chunk)
            print(f"Summarizing chunk {i+1}/{len(chunks)} ({chunk_tokens} tokens)...")
            max_len = min(chunk_tokens // 2, 400)
            summary = self._summarize_single(chunk, max_length=max_len, min_length=40)
            chunk_summaries.append(summary)

        #   c. combined = join all chunk_summaries together
        combined = " ".join(chunk_summaries)
        combined_tokens = self._count_tokens(combined)
        print(f"Combined chunk summaries into one text ({combined_tokens} tokens).")

        #   d. final_summary = self._summarize_single(combined)  ← second pass
        if combined_tokens > self.max_input_tokens:
            print(f"Combined summary still too long ({combined_tokens} tokens), summarizing again.")
            return self.summarize(combined, final_pass=final_pass)
        else:
            if final_pass:
                print(f"    Final pass: generating coherent summary...")
                max_len = min(combined_tokens, 800)
                return self._summarize_single(combined, max_length=max_len, min_length=50)
            else:
                print(f"    Returning concatenated summaries ({len(combined.split())} words)")
                return combined


# --- Long-T5 Summarizer ---
class LongT5Summarizer(BARTSummarizer):
    """Long-T5 model for long-document summarization. Inherits all logic from BARTSummarizer."""

    def __init__(self):
        self.model_name = "pszemraj/long-t5-tglobal-base-16384-book-summary"
        self.pipe = pipeline("summarization", model=self.model_name, device=DEVICE)
        self.tokenizer = AutoTokenizer.from_pretrained(self.model_name)
        self.max_input_tokens = 16384


# --- BARTSamsumSummarizer ---
class BARTSamsumSummarizer(BARTSummarizer):
    """BART fine-tuned on SAMSum dialogue summarization dataset."""

    def __init__(self):
        self.model_name = "philschmid/bart-large-cnn-samsum"
        self.pipe = pipeline("summarization", model=self.model_name, device=DEVICE)
        self.tokenizer = AutoTokenizer.from_pretrained(self.model_name)
        self.max_input_tokens = 1024



# --- LED Arxiv Summarizer ---
class LEDArxivSummarizer(BARTSummarizer):
    """Longformer Encoder-Decoder fine-tuned on arXiv papers."""

    def __init__(self):
        self.model_name = "allenai/led-large-16384-arxiv"
        self.pipe = pipeline("summarization", model=self.model_name, device=DEVICE)
        self.tokenizer = AutoTokenizer.from_pretrained(self.model_name)
        self.max_input_tokens = 16384


# --- QwenSummarizer class ---
class QwenSummarizer(BaseSummarizer):
    """
    LLM-based summarizer using Qwen2.5-7B-Instruct with 4-bit quantization.
    Unlike BART/T5 models, this uses text-generation pipeline with prompting.
    128K context window allows single-pass summarization without chunking.
    """

    def __init__(self, output_ratio=0.06):
        from transformers import AutoModelForCausalLM, AutoTokenizer, BitsAndBytesConfig

        self.model_name = "Qwen/Qwen2.5-7B-Instruct"

        # 4-bit quantization config to fit in GPU memory
        bnb_config = BitsAndBytesConfig(
            load_in_4bit=True,
            bnb_4bit_compute_dtype=torch.float16,  # compute in float16 for speed
            bnb_4bit_quant_type="nf4",              # normalized float 4-bit
        )

        # Load tokenizer
        self.tokenizer = AutoTokenizer.from_pretrained(self.model_name)

        # Load model with quantization
        self.model = AutoModelForCausalLM.from_pretrained(
            self.model_name,
            quantization_config=bnb_config,
            device_map="auto",   # automatically place layers on available GPU(s)
        )

        # TODO: to see whether we should modify this value
        self.max_input_tokens = 65536  # use 65K of the 128K window, leave room for output

        self.output_ratio = output_ratio

    def _count_tokens(self, text):
        return len(self.tokenizer.encode(text, add_special_tokens=False))


    def summarize(self, text: str, final_pass: bool = True) -> str:
        """
        Summarize using LLM with prompting.
        final_pass parameter is accepted for compatibility with process_all_topics
        but has no effect — LLM always produces a coherent final summary.
        """
        total_tokens = self._count_tokens(text)
        print(f"Input text: {total_tokens} tokens.")

        # If input exceeds our limit, truncate from the end
        # (lectures often have Q&A / chit-chat at the end which is less important)
        if total_tokens > self.max_input_tokens:
            print(f"Input too long ({total_tokens} tokens), truncating to {self.max_input_tokens}.")
            tokens = self.tokenizer.encode(text, add_special_tokens=False)
            text = self.tokenizer.decode(tokens[:self.max_input_tokens])

        # Build prompt using Qwen's chat template
        messages = [
            {"role": "system", "content": (
                "You are an expert academic summarizer specializing in computer "
                "science and deep learning courses. You produce clear, "
                "well-organized summaries from university lecture materials."
            )},
            {"role": "user", "content": (
                "Summarize the following university lecture material into a comprehensive "
                "yet concise summary. The input combines slide text, lecture transcript, and student "
                "notes, so it contains filler words, informal language, typos, and "
                "off-topic conversations — ignore all of these.\n\n"
                "Your summary should:\n"
                "- Cover every major topic and subtopic discussed\n"
                "- Include key definitions, formulas, and technical details\n"
                "- Preserve the logical flow of the lecture\n"
                "- Use clear academic language\n\n"
                f"INPUT TEXT:\n{text}"
            )},
        ]

        # Apply chat template (Qwen uses a specific format for instruction following)
        prompt = self.tokenizer.apply_chat_template(
            messages,
            tokenize=False,
            add_generation_prompt=True,
        )

        # Tokenize and generate
        inputs = self.tokenizer(prompt, return_tensors="pt").to("cuda")

        # Dynamic output length
        target_output = int(total_tokens * self.output_ratio)  # output ratio of input
        max_new = max(300, target_output)
        min_new = max(150, max_new // 4)

        print("Generating summary...")
        with torch.no_grad():
            outputs = self.model.generate(
                **inputs,
                max_new_tokens=max_new,     # max output length (~500 words)
                min_new_tokens=min_new,     # min output length (~150 words)
                do_sample=False,        # greedy decoding for reproducibility
                temperature=1.0,
                repetition_penalty=1.1, # slight penalty to avoid repetitive output
            )

        # Decode only the NEW tokens (exclude the prompt)
        generated_tokens = outputs[0][inputs["input_ids"].shape[1]:]
        summary = self.tokenizer.decode(generated_tokens, skip_special_tokens=True)

        print(f"Generated summary: {len(summary.split())} words.")
        return summary.strip()


# --- Claude Sonnet Summarizer (API-based) ---
class ClaudeSummarizer(BaseSummarizer):
    """Claude Sonnet 4 via Anthropic API. No chunking needed — 200K context."""

    def __init__(self):
        import anthropic
        api_key = os.environ.get("ANTHROPIC_API_KEY")
        if not api_key or api_key == "paste-your-key-here":
            raise ValueError(
                "Set ANTHROPIC_API_KEY in .env file. "
                "Get one at https://console.anthropic.com"
            )
        self.client = anthropic.Anthropic(api_key=api_key)
        self.model_name = "claude-sonnet-4-20250514"
        self.max_input_tokens = 200000  # 200K context window

    def summarize(self, text: str, final_pass: bool = True) -> str:
        word_count = len(text.split())
        print(f"Sending {word_count} words to {self.model_name} ...")

        if final_pass:
            # Final strategy: concise, coherent summary
            prompt = (
                "You are an expert educational summarizer. "
                "Summarize the following lecture content into a concise, coherent summary "
                "suitable for a student reviewing for exams. "
                "Focus on key concepts, definitions, methods, and relationships. "
                "Be factual — only include information present in the source material. "
                "Keep the summary between 100-300 words.\n\n"
                f"LECTURE CONTENT:\n{text}"
            )
        else:
            # Concat strategy: detailed, structured summary
            prompt = (
                "You are an expert educational summarizer. "
                "Create a detailed, structured summary of the following lecture content. "
                "Cover all major topics and subtopics with key details. "
                "Use clear section headers. Be factual — only include information "
                "present in the source material. "
                "Aim for 300-600 words.\n\n"
                f"LECTURE CONTENT:\n{text}"
            )

        response = self.client.messages.create(
            model=self.model_name,
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}],
        )
        summary = response.content[0].text
        print(f"  Received {len(summary.split())} words from API.")
        return summary


# --- Helper function to process all files ---
def process_all_topics(input_dir, output_dir, model, model_tag="default", strategy="final"):
    """
    Loop through all _ori.txt files in input_dir,
    summarize each, save as _sum.txt in output_dir.
    """
    input_path = Path(input_dir)
    output_path = Path(output_dir) / model_tag / strategy
    output_path.mkdir(parents=True, exist_ok=True)

    # Step 1: Find all files ending with _ori.txt in input_dir
    ori_files = sorted(input_path.glob("*_ori.txt"))

    if not ori_files:
        print(f"No files ending with '_ori.txt' found in {input_dir}.")
        return {}
    print(f"Found {len(ori_files)} files to process.")

    results = {}

    # Step 2: For each file:
    for file in ori_files:
    #     a. Read the text
        text = file.read_text(encoding="utf-8")
        print(f"Input length: {len(text.split())} words")

    #     b. summary = model.summarize(text)
        final_pass = (strategy == "final")
        summary = model.summarize(text, final_pass=final_pass)
        print(f"Summary length: {len(summary.split())} words")

        # Create output filename: use only the last part of model_tag
        short_tag = model_tag.split("/")[-1]
        output_filename = file.name.replace("_ori.txt", f"_sum_{short_tag}_{strategy}.txt")
        output_file = output_path / output_filename

    #     d. Write summary to output_dir / output_filename
        output_file.write_text(summary, encoding="utf-8")
        print(f"Saved summary to {output_file}")

    #     e. Print progress: which file, token count, summary length
        results[file.name] = {
            "input_file": str(file),
            "output_file": str(output_file),
            "input_words": len(text.split()),
            "summary_words": len(summary.split()),
        }

    print(f"Summarized {len(results)} files.")
        
    return results



# --- Test block ---
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="StudyLens Summarization Pipeline")
    parser.add_argument("--model", type=str, required=True,
                        choices=["first5", "random", "tfidf", "bart", "longt5", "bart-samsum", "led-arxiv", "qwen7b"],
                        help="Which model to run")
    parser.add_argument("--num_sentences", type=int, default=15,
                        help="Number of sentences for extractive models (default: 15)")
    parser.add_argument("--strategy", type=str, default="both",
                        choices=["concat", "final", "both"],
                        help="Summarization strategy (default: both)")
    parser.add_argument("--input_dir", type=str, default="data/processed")
    parser.add_argument("--output_dir", type=str, default="data/outputs")
    parser.add_argument("--output_ratio", type=float, default=0.06,
                        help="Output length as fraction of input tokens (default: 0.06)")
    args = parser.parse_args()

    # Model dispatch
    if args.model == "first5":
        model = FirstSentenceSummarizer(num_slides=5)
    elif args.model == "random":
        model = RandomExtractiveSummarizer(num_sentences=args.num_sentences)
    elif args.model == "tfidf":
        model = TFIDFExtractiveSummarizer(num_sentences=args.num_sentences)
    elif args.model == "bart":
        model = BARTSummarizer()
    elif args.model == "longt5":
        model = LongT5Summarizer()
    elif args.model == "bart-samsum":
        model = BARTSamsumSummarizer()
    elif args.model == "led-arxiv":
        model = LEDArxivSummarizer()
    elif args.model == "qwen7b":
        model = QwenSummarizer(output_ratio=args.output_ratio)

    # Run strategies
    if args.model == "first5":
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag="naive", strategy="first5")
    elif args.model == "random":
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag="naive", strategy="random")
    elif args.model == "tfidf":
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag="classical_ml", strategy="tfidf")
    elif args.model == "qwen7b":
        ratio_tag = f"ratio{int(args.output_ratio * 100):02d}"
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag="neural_network/qwen7b", strategy=ratio_tag)
    elif args.strategy == "both":
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag=f"neural_network/{args.model}", strategy="concat")
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag=f"neural_network/{args.model}", strategy="final")
    else:
        process_all_topics(args.input_dir, args.output_dir, model,
                           model_tag=f"neural_network/{args.model}", strategy=args.strategy)

    print("\nDone!")